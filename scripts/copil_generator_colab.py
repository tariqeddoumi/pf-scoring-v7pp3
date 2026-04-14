#!/usr/bin/env python3
"""
COPIL Presentation Generator - Google Colab Version
Generates professional PowerPoint presentations from CSV data
"""

import csv
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Colors for professional banking theme
PRIMARY_BLUE = RGBColor(0, 51, 102)      # #003366
ACCENT_ORANGE = RGBColor(255, 153, 0)   # #FF9900
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(76, 175, 80)
ORANGE = RGBColor(255, 152, 0)
RED = RGBColor(244, 67, 54)
DARK_TEXT = RGBColor(33, 33, 33)


def parse_completion(item):
    """Safely parse completion percentage"""
    try:
        comp_str = item.get("COMPLÉTION_%", "0").rstrip("%").strip()
        return float(comp_str) if comp_str else 0
    except:
        return 0


def load_csv_data(csv_file):
    """Load CSV data into list of dictionaries"""
    print(f"📥 Loading CSV data from {csv_file}...")

    if not os.path.exists(csv_file):
        print(f"❌ Error: {csv_file} not found")
        return []

    items = []
    with open(csv_file, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        items = list(reader)

    print(f"✅ Loaded {len(items)} items from CSV")
    return items


def calculate_statistics(items):
    """Calculate KPIs and statistics"""
    print("📊 Calculating statistics...")

    # Parse completion percentage and calculate average
    completions = [parse_completion(item) for item in items]

    stats = {
        "total": len(items),
        "integrated": len([i for i in items if i.get("STATUT") == "INTÉGRÉ"]),
        "in_progress": len([i for i in items if i.get("STATUT") == "EN COURS"]),
        "planned": len([i for i in items if i.get("STATUT") == "PLANIFIÉ"]),
        "blocked": len([i for i in items if i.get("STATUT") == "BLOQUÉ"]),
        "by_bloc": {},
        "risks": [i for i in items if i.get("STATUT") == "BLOQUÉ" or parse_completion(i) < 50],
    }

    stats["average_completion"] = sum(completions) / len(completions) if completions else 0

    # Group by bloc
    for item in items:
        bloc = item.get("BLOC", "Unknown")
        if bloc not in stats["by_bloc"]:
            stats["by_bloc"][bloc] = {
                "total": 0,
                "integrated": 0,
                "items": [],
                "completion_sum": 0,
            }

        stats["by_bloc"][bloc]["total"] += 1
        if item.get("STATUT") == "INTÉGRÉ":
            stats["by_bloc"][bloc]["integrated"] += 1
        stats["by_bloc"][bloc]["items"].append(item)
        stats["by_bloc"][bloc]["completion_sum"] += parse_completion(item)

    # Calculate completion percentage for each bloc
    for bloc in stats["by_bloc"]:
        items_count = stats["by_bloc"][bloc]["total"]
        stats["by_bloc"][bloc]["completion"] = (
            stats["by_bloc"][bloc]["completion_sum"] / items_count
            if items_count > 0
            else 0
        )

    print(f"✅ Statistics calculated:")
    print(f"   - Total items: {stats['total']}")
    print(f"   - Integrated: {stats['integrated']}")
    print(f"   - In progress: {stats['in_progress']}")
    print(f"   - Average completion: {stats['average_completion']:.1f}%")

    return stats


def add_title_slide(prs, stats):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "PF SCORING"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "COMITÉ DE PILOTAGE - COPIL"
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Statut Global: {stats['average_completion']:.0f}% Complété"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_kpi_slide(prs, stats):
    """Add KPI summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "RÉSUMÉ EXÉCUTIF - KPIs"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # KPI Cards
    kpis = [
        ("INTÉGRÉ", stats["integrated"], GREEN),
        ("EN COURS", stats["in_progress"], ORANGE),
        ("PLANIFIÉ", stats["planned"], RGBColor(33, 150, 243)),
        ("BLOQUÉ", stats["blocked"], RED),
    ]

    card_width = 1.8
    card_height = 1.5
    start_x = 0.8
    start_y = 1.3

    for idx, (label, count, color) in enumerate(kpis):
        x = start_x + (idx * 2.1)
        # Card background
        shape = slide.shapes.add_shape(1, Inches(x), Inches(start_y), Inches(card_width), Inches(card_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color

        # Count text
        count_box = slide.shapes.add_textbox(Inches(x), Inches(start_y + 0.2), Inches(card_width), Inches(0.7))
        count_frame = count_box.text_frame
        p = count_frame.paragraphs[0]
        p.text = str(count)
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label text
        label_box = slide.shapes.add_textbox(Inches(x), Inches(start_y + 0.9), Inches(card_width), Inches(0.4))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Summary text
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True
    p = summary_frame.paragraphs[0]
    p.text = f"Total des éléments: {stats['total']} | Complétude moyenne: {stats['average_completion']:.1f}%"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER


def add_completion_chart_slide(prs, stats):
    """Add completion chart by bloc"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "AVANCEMENT PAR BLOC"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = list(stats["by_bloc"].keys())

    completions = [stats["by_bloc"][bloc]["completion"] for bloc in stats["by_bloc"]]
    chart_data.add_series("Complétude %", tuple(completions))

    # Add chart
    x, y, cx, cy = Inches(0.5), Inches(1.1), Inches(9), Inches(5.2)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy).chart

    chart.has_legend = False
    chart_data = chart.chart_data
    chart_data.categories = list(stats["by_bloc"].keys())
    chart_data.add_series("Complétude %", tuple(completions))

    # Format chart
    plot = chart.plots[0]
    plot.vary_by_categories = True
    series = plot.series[0]

    # Color series based on completion
    for idx, (bloc, completion) in enumerate(
        zip(stats["by_bloc"].keys(), completions)
    ):
        if completion >= 90:
            color = GREEN
        elif completion >= 70:
            color = ORANGE
        else:
            color = RED
        series.points[idx].format.fill.solid()
        series.points[idx].format.fill.fore_color.rgb = color


def add_gauge_slide(prs, stats):
    """Add global completion gauge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "COMPLÉTUDE GLOBALE"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Large percentage display
    completion = stats["average_completion"]
    percent_box = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(4), Inches(2))
    percent_frame = percent_box.text_frame
    p = percent_frame.paragraphs[0]
    p.text = f"{completion:.0f}%"
    p.font.size = Pt(120)
    p.font.bold = True
    if completion >= 80:
        p.font.color.rgb = GREEN
    elif completion >= 50:
        p.font.color.rgb = ORANGE
    else:
        p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    # Progress bar background
    bar_bg = slide.shapes.add_shape(1, Inches(1.5), Inches(4.8), Inches(7), Inches(0.5))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = LIGHT_GRAY
    bar_bg.line.color.rgb = RGBColor(200, 200, 200)

    # Progress bar foreground
    bar_width = 7 * (completion / 100)
    bar_fg = slide.shapes.add_shape(1, Inches(1.5), Inches(4.8), Inches(bar_width), Inches(0.5))
    bar_fg.fill.solid()
    if completion >= 80:
        bar_fg.fill.fore_color.rgb = GREEN
    elif completion >= 50:
        bar_fg.fill.fore_color.rgb = ORANGE
    else:
        bar_fg.fill.fore_color.rgb = RED
    bar_fg.line.color.rgb = bar_fg.fill.fore_color.rgb

    # Status text
    status_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.8))
    status_frame = status_box.text_frame
    p = status_frame.paragraphs[0]
    if completion >= 80:
        status = "✅ Avancement Excellent"
    elif completion >= 50:
        status = "⚠️  Avancement Acceptable"
    else:
        status = "❌ Avancement Faible"
    p.text = status
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER


def add_risks_slide(prs, stats):
    """Add risks and blockers slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "RISQUES & BLOCAGES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    risks = stats["risks"][:5]  # Top 5 risks

    if not risks:
        no_risks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        no_risks_frame = no_risks_box.text_frame
        p = no_risks_frame.paragraphs[0]
        p.text = "✅ Aucun risque majeur identifié"
        p.font.size = Pt(28)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
    else:
        y = 1.3
        for idx, risk in enumerate(risks):
            # Risk item
            risk_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.8))
            risk_frame = risk_box.text_frame
            risk_frame.word_wrap = True
            p = risk_frame.paragraphs[0]
            element = risk.get("ÉLÉMENT", "Unknown")
            status = risk.get("STATUT", "PLANIFIÉ")
            completion = parse_completion(risk)
            p.text = f"🔴 {element} ({status} - {completion:.0f}%)"
            p.font.size = Pt(14)
            p.font.color.rgb = RED
            p.space_after = Pt(6)

            y += 0.9


def add_next_steps_slide(prs, stats):
    """Add next steps slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "PROCHAINES ÉTAPES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    phases = [
        ("Phase 9", "Design Responsive Mobile"),
        ("Phase 10", "Tests & QA Complète"),
        ("Phase 11", "Déploiement Production"),
        ("Phase 12", "Monitoring & Support"),
    ]

    y = 1.3
    for phase, description in phases:
        # Phase box
        phase_shape = slide.shapes.add_shape(
            1, Inches(0.7), Inches(y), Inches(0.6), Inches(0.6)
        )
        phase_shape.fill.solid()
        phase_shape.fill.fore_color.rgb = ACCENT_ORANGE
        phase_shape.line.color.rgb = ACCENT_ORANGE

        # Phase number
        phase_text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(0.6), Inches(0.6))
        phase_text_frame = phase_text_box.text_frame
        phase_text_frame.vertical_anchor = 1  # Middle
        p = phase_text_frame.paragraphs[0]
        p.text = phase.split()[1]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.05), Inches(7.8), Inches(0.5))
        desc_frame = desc_box.text_frame
        p = desc_frame.paragraphs[0]
        p.text = f"{phase} - {description}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT

        y += 1


def create_presentation(stats, output_file="COPIL_Presentation.pptx"):
    """Create and save presentation"""
    print("\n📝 Creating presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    add_title_slide(prs, stats)
    add_kpi_slide(prs, stats)
    add_completion_chart_slide(prs, stats)
    add_gauge_slide(prs, stats)
    add_risks_slide(prs, stats)
    add_next_steps_slide(prs, stats)

    # Save
    prs.save(output_file)
    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_file


def main(csv_file="PF_SCORING_SPECIFICATIONS_TRACKING.csv", output_file="COPIL_Presentation.pptx"):
    """Main function"""
    print("\n🚀 PF Scoring - COPIL Presentation Generator (Colab Edition)")
    print("=" * 60)

    # Load data
    items = load_csv_data(csv_file)
    if not items:
        print("❌ No data to process")
        return

    # Calculate statistics
    stats = calculate_statistics(items)

    # Create presentation
    create_presentation(stats, output_file)

    print("\n" + "=" * 60)
    print("✅ Presentation ready!")
    print(f"📥 Download: {output_file}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    import sys
    csv_file = sys.argv[1] if len(sys.argv) > 1 else "PF_SCORING_SPECIFICATIONS_TRACKING.csv"
    output_file = sys.argv[2] if len(sys.argv) > 2 else "COPIL_Presentation.pptx"
    main(csv_file, output_file)
