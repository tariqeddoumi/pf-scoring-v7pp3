#!/usr/bin/env python3
"""
COPIL Presentation Generator - Google Colab Version v4
Generates professional PowerPoint presentations from API data
Integrates with PF Scoring backend API
"""

import requests
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Colors for professional banking theme
PRIMARY_BLUE = RGBColor(0, 51, 102)      # #003366
ACCENT_ORANGE = RGBColor(255, 153, 0)   # #FF9900
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(76, 175, 80)
ORANGE = RGBColor(255, 152, 0)
RED = RGBColor(244, 67, 54)
DARK_TEXT = RGBColor(33, 33, 33)
LIGHT_BLUE = RGBColor(33, 150, 243)


class PFScoringAPIClient:
    """Client for PF Scoring API"""

    def __init__(self, api_url, auth_token=None):
        """Initialize API client"""
        self.api_url = api_url.rstrip("/")
        self.auth_token = auth_token
        self.headers = {
            "Content-Type": "application/json"
        }
        if auth_token:
            self.headers["Authorization"] = f"Bearer {auth_token}"

    def _request(self, method, endpoint, **kwargs):
        """Make API request"""
        url = f"{self.api_url}/api/{endpoint}"
        try:
            response = requests.request(method, url, headers=self.headers, **kwargs)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.RequestException as e:
            print(f"❌ API Error: {e}")
            return None

    def get_users(self):
        """Fetch all users"""
        return self._request("GET", "users") or []

    def get_projects(self):
        """Fetch all projects"""
        return self._request("GET", "projects") or []

    def get_evaluations(self):
        """Fetch all evaluations"""
        return self._request("GET", "evaluations") or []

    def get_stats(self):
        """Fetch statistics"""
        return self._request("GET", "stats") or {}


def collect_data_from_api(api_url, auth_token=None):
    """Collect data from API"""
    print(f"📡 Connecting to API: {api_url}...")

    client = PFScoringAPIClient(api_url, auth_token)

    print("📥 Fetching users...")
    users_response = client.get_users()
    users = users_response.get("data", []) if isinstance(users_response, dict) else users_response
    print(f"✅ Loaded {len(users)} users")

    print("📥 Fetching projects...")
    projects_response = client.get_projects()
    projects = projects_response.get("data", []) if isinstance(projects_response, dict) else projects_response
    print(f"✅ Loaded {len(projects)} projects")

    print("📥 Fetching evaluations...")
    evaluations_response = client.get_evaluations()
    evaluations = evaluations_response.get("data", []) if isinstance(evaluations_response, dict) else evaluations_response
    print(f"✅ Loaded {len(evaluations)} evaluations")

    return {
        "users": users,
        "projects": projects,
        "evaluations": evaluations
    }


def collect_data_from_csv(csv_file):
    """Fallback: Collect data from CSV file for COPIL tracking"""
    print(f"📥 Loading CSV data from {csv_file}...")

    import csv

    if not Path(csv_file).exists():
        print(f"❌ Error: {csv_file} not found")
        return {"items": []}

    items = []
    with open(csv_file, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        items = list(reader)

    print(f"✅ Loaded {len(items)} items from CSV")
    return {"items": items}


def calculate_statistics_from_api_data(data):
    """Calculate KPIs from API data"""
    print("📊 Calculating statistics from API data...")

    users = data.get("users", [])
    projects = data.get("projects", [])
    evaluations = data.get("evaluations", [])

    # User statistics
    stats = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "total_users": len(users),
        "users_by_role": {},
        "total_projects": len(projects),
        "projects_by_status": {},
        "total_evaluations": len(evaluations),
        "evaluations_by_status": {},
        "evaluation_workflow": {
            "draft": 0,
            "submitted": 0,
            "validated": 0,
            "approved": 0,
            "rejected": 0
        }
    }

    # Count users by role
    for user in users:
        role = user.get("role", "guest").lower()
        if role not in stats["users_by_role"]:
            stats["users_by_role"][role] = 0
        stats["users_by_role"][role] += 1

    # Count projects by status
    for project in projects:
        status = project.get("status", "draft").lower()
        if status not in stats["projects_by_status"]:
            stats["projects_by_status"][status] = 0
        stats["projects_by_status"][status] += 1

    # Count evaluations by status
    for eval in evaluations:
        status = eval.get("status", "draft").lower()
        if status not in stats["evaluations_by_status"]:
            stats["evaluations_by_status"][status] = 0
        stats["evaluations_by_status"][status] += 1

        # Map to workflow stages
        if status == "draft":
            stats["evaluation_workflow"]["draft"] += 1
        elif status == "submitted":
            stats["evaluation_workflow"]["submitted"] += 1
        elif status == "validated":
            stats["evaluation_workflow"]["validated"] += 1
        elif status == "approved":
            stats["evaluation_workflow"]["approved"] += 1
        elif status == "rejected":
            stats["evaluation_workflow"]["rejected"] += 1

    print(f"✅ Statistics calculated:")
    print(f"   - Total users: {stats['total_users']}")
    print(f"   - Total projects: {stats['total_projects']}")
    print(f"   - Total evaluations: {stats['total_evaluations']}")

    return stats


def calculate_statistics_from_csv(data):
    """Calculate KPIs from CSV data (legacy COPIL tracking)"""
    print("📊 Calculating statistics from CSV data...")

    items = data.get("items", [])

    def parse_completion(item):
        """Safely parse completion percentage"""
        try:
            comp_str = item.get("COMPLÉTION_%", "0").rstrip("%").strip()
            return float(comp_str) if comp_str else 0
        except:
            return 0

    stats = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "total": len(items),
        "integrated": len([i for i in items if i.get("STATUT") == "INTÉGRÉ"]),
        "in_progress": len([i for i in items if i.get("STATUT") == "EN COURS"]),
        "planned": len([i for i in items if i.get("STATUT") == "PLANIFIÉ"]),
        "blocked": len([i for i in items if i.get("STATUT") == "BLOQUÉ"]),
        "by_bloc": {},
        "risks": [i for i in items if i.get("STATUT") == "BLOQUÉ" or parse_completion(i) < 50],
    }

    completions = [parse_completion(item) for item in items]
    stats["average_completion"] = sum(completions) / len(completions) if completions else 0

    # Group by bloc
    for item in items:
        bloc = item.get("BLOC", "Unknown")
        if bloc not in stats["by_bloc"]:
            stats["by_bloc"][bloc] = {
                "total": 0,
                "integrated": 0,
                "items": [],
                "completion_sum": 0,
            }

        stats["by_bloc"][bloc]["total"] += 1
        if item.get("STATUT") == "INTÉGRÉ":
            stats["by_bloc"][bloc]["integrated"] += 1
        stats["by_bloc"][bloc]["items"].append(item)
        stats["by_bloc"][bloc]["completion_sum"] += parse_completion(item)

    for bloc in stats["by_bloc"]:
        items_count = stats["by_bloc"][bloc]["total"]
        stats["by_bloc"][bloc]["completion"] = (
            stats["by_bloc"][bloc]["completion_sum"] / items_count
            if items_count > 0
            else 0
        )

    print(f"✅ Statistics calculated:")
    print(f"   - Total items: {stats['total']}")
    print(f"   - Integrated: {stats['integrated']}")

    return stats


def add_title_slide(prs, stats):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "PF SCORING"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "COMITÉ DE PILOTAGE - COPIL"
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Timestamp
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Rapport généré: {stats.get('timestamp', 'N/A')}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_kpi_users_slide(prs, stats):
    """Add KPI slide for users and roles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "GESTION DES UTILISATEURS"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Total users
    total_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.6))
    total_frame = total_box.text_frame
    p = total_frame.paragraphs[0]
    p.text = f"Total Utilisateurs: {stats.get('total_users', 0)}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Users by role
    users_by_role = stats.get("users_by_role", {})
    y = 2.0

    role_colors = {
        "admin": RGBColor(244, 67, 54),
        "manager": RGBColor(255, 152, 0),
        "analyst": RGBColor(76, 175, 80),
        "guest": RGBColor(200, 200, 200),
        "viewer": LIGHT_BLUE
    }

    for role, count in sorted(users_by_role.items()):
        # Role name
        role_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(3), Inches(0.4))
        role_frame = role_box.text_frame
        p = role_frame.paragraphs[0]
        p.text = f"{role.upper()}: {count}"
        p.font.size = Pt(16)
        p.font.color.rgb = role_colors.get(role.lower(), DARK_TEXT)

        # Progress bar
        bar_width = (count / max(users_by_role.values(), default=1)) * 4 if users_by_role else 0
        bar = slide.shapes.add_shape(1, Inches(4), Inches(y + 0.05), Inches(bar_width), Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = role_colors.get(role.lower(), DARK_TEXT)
        bar.line.color.rgb = role_colors.get(role.lower(), DARK_TEXT)

        y += 0.6


def add_kpi_projects_slide(prs, stats):
    """Add KPI slide for projects"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "GESTION DES PROJETS"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Total projects
    total_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.6))
    total_frame = total_box.text_frame
    p = total_frame.paragraphs[0]
    p.text = f"Total Projets: {stats.get('total_projects', 0)}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Projects by status
    projects_by_status = stats.get("projects_by_status", {})

    status_colors = {
        "draft": RGBColor(200, 200, 200),
        "in_review": ORANGE,
        "approved": GREEN,
        "rejected": RED
    }

    y = 2.0
    for status, count in sorted(projects_by_status.items()):
        # Status name
        status_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(3), Inches(0.4))
        status_frame = status_box.text_frame
        p = status_frame.paragraphs[0]
        p.text = f"{status.upper()}: {count}"
        p.font.size = Pt(16)
        p.font.color.rgb = status_colors.get(status.lower(), DARK_TEXT)

        # Progress bar
        bar_width = (count / max(projects_by_status.values(), default=1)) * 4 if projects_by_status else 0
        bar = slide.shapes.add_shape(1, Inches(4), Inches(y + 0.05), Inches(bar_width), Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = status_colors.get(status.lower(), DARK_TEXT)
        bar.line.color.rgb = status_colors.get(status.lower(), DARK_TEXT)

        y += 0.6


def add_evaluation_workflow_slide(prs, stats):
    """Add slide for evaluation workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "FLUX DES ÉVALUATIONS"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Total evaluations
    total_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    total_frame = total_box.text_frame
    p = total_frame.paragraphs[0]
    p.text = f"Total Évaluations: {stats.get('total_evaluations', 0)}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Workflow stages
    workflow = stats.get("evaluation_workflow", {})
    stages = ["draft", "submitted", "validated", "approved", "rejected"]
    stage_labels = ["Brouillon", "Soumise", "Validée", "Approuvée", "Rejetée"]

    y = 2.0
    for stage, label in zip(stages, stage_labels):
        count = workflow.get(stage, 0)
        percentage = (count / stats.get('total_evaluations', 1) * 100) if stats.get('total_evaluations') else 0

        # Stage name and count
        stage_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(3), Inches(0.4))
        stage_frame = stage_box.text_frame
        p = stage_frame.paragraphs[0]
        p.text = f"{label}: {count} ({percentage:.0f}%)"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT

        # Progress bar
        bar_width = (count / max(workflow.values(), default=1)) * 4 if workflow else 0
        bar = slide.shapes.add_shape(1, Inches(4), Inches(y + 0.05), Inches(bar_width), Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT_BLUE
        bar.line.color.rgb = PRIMARY_BLUE

        y += 0.6


def add_next_steps_slide(prs):
    """Add next steps slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "PROCHAINES ÉTAPES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    phases = [
        ("Phase 9", "Gestion d'Accès & Authentification"),
        ("Phase 10", "Tests & Validation Complète"),
        ("Phase 11", "Déploiement Production"),
        ("Phase 12", "Monitoring & Support"),
    ]

    y = 1.3
    for phase, description in phases:
        # Phase box
        phase_shape = slide.shapes.add_shape(
            1, Inches(0.7), Inches(y), Inches(0.6), Inches(0.6)
        )
        phase_shape.fill.solid()
        phase_shape.fill.fore_color.rgb = ACCENT_ORANGE
        phase_shape.line.color.rgb = ACCENT_ORANGE

        # Phase number
        phase_text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(0.6), Inches(0.6))
        phase_text_frame = phase_text_box.text_frame
        phase_text_frame.vertical_anchor = 1  # Middle
        p = phase_text_frame.paragraphs[0]
        p.text = phase.split()[1]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.05), Inches(7.8), Inches(0.5))
        desc_frame = desc_box.text_frame
        p = desc_frame.paragraphs[0]
        p.text = f"{phase} - {description}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT

        y += 1


def create_presentation(stats, is_api_data=True, output_file="COPIL_Presentation.pptx"):
    """Create and save presentation"""
    print("\n📝 Creating presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    add_title_slide(prs, stats)

    if is_api_data:
        add_kpi_users_slide(prs, stats)
        add_kpi_projects_slide(prs, stats)
        add_evaluation_workflow_slide(prs, stats)

    add_next_steps_slide(prs)

    # Save
    prs.save(output_file)
    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_file


def main(api_url=None, auth_token=None, csv_file=None, output_file="COPIL_Presentation.pptx"):
    """Main function"""
    print("\n🚀 PF Scoring - COPIL Presentation Generator v4 (Colab Edition)")
    print("=" * 60)

    # Collect data
    if api_url:
        data = collect_data_from_api(api_url, auth_token)
        stats = calculate_statistics_from_api_data(data)
        is_api = True
    elif csv_file:
        data = collect_data_from_csv(csv_file)
        stats = calculate_statistics_from_csv(data)
        is_api = False
    else:
        print("❌ Either API URL or CSV file required")
        return

    # Create presentation
    create_presentation(stats, is_api, output_file)

    print("\n" + "=" * 60)
    print("✅ Presentation ready!")
    print(f"📥 Download: {output_file}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    import sys

    # Parse arguments
    api_url = None
    auth_token = None
    csv_file = None
    output_file = "COPIL_Presentation.pptx"

    # Check for command line arguments
    if len(sys.argv) > 1:
        api_url = sys.argv[1] if sys.argv[1].startswith("http") else None
        csv_file = sys.argv[1] if not sys.argv[1].startswith("http") else None

    if len(sys.argv) > 2:
        auth_token = sys.argv[2]

    if len(sys.argv) > 3:
        output_file = sys.argv[3]

    main(api_url, auth_token, csv_file, output_file)
