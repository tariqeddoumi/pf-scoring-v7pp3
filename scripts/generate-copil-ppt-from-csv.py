#!/usr/bin/env python3
"""
Generate COPIL Presentation PPT from CSV data
Creates professional PowerPoint presentations for project governance
"""

import csv
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color scheme (Banque Marocaine Professional)
COLORS = {
    "primary": RGBColor(0, 51, 102),  # Dark blue
    "secondary": RGBColor(0, 102, 153),  # Medium blue
    "accent": RGBColor(255, 153, 0),  # Orange
    "green": RGBColor(76, 175, 80),  # Green
    "orange": RGBColor(255, 152, 0),  # Orange
    "red": RGBColor(244, 67, 54),  # Red
    "text": RGBColor(33, 33, 33),  # Dark gray
    "light_text": RGBColor(117, 117, 117),  # Light gray
    "white": RGBColor(255, 255, 255),  # White
}


def load_csv_data(csv_file="PF_SCORING_SPECIFICATIONS_TRACKING.csv"):
    """Load tracking data from CSV"""
    print("📥 Loading CSV data...")

    items = []
    with open(csv_file, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            items.append(row)

    print(f"✅ Loaded {len(items)} items from CSV")
    return items


def parse_completion(item):
    """Safely parse completion percentage"""
    try:
        comp_str = item.get("COMPLÉTION_%", "0").rstrip("%").strip()
        return float(comp_str) if comp_str else 0
    except:
        return 0


def calculate_statistics(items):
    """Calculate KPIs and statistics"""
    print("📊 Calculating statistics...")

    # Parse completion percentage and calculate average
    completions = [parse_completion(item) for item in items]

    stats = {
        "total": len(items),
        "integrated": len([i for i in items if i.get("STATUT") == "INTÉGRÉ"]),
        "in_progress": len([i for i in items if i.get("STATUT") == "EN COURS"]),
        "planned": len([i for i in items if i.get("STATUT") == "PLANIFIÉ"]),
        "blocked": len([i for i in items if i.get("STATUT") == "BLOQUÉ"]),
        "by_bloc": {},
        "risks": [i for i in items if i.get("STATUT") == "BLOQUÉ" or parse_completion(i) < 50],
    }

    stats["average_completion"] = sum(completions) / len(completions) if completions else 0

    # Group by bloc
    for item in items:
        bloc = item.get("BLOC", "Unknown")
        if bloc not in stats["by_bloc"]:
            stats["by_bloc"][bloc] = {
                "total": 0,
                "integrated": 0,
                "items": [],
                "completion_sum": 0,
            }

        stats["by_bloc"][bloc]["total"] += 1
        if item.get("STATUT") == "INTÉGRÉ":
            stats["by_bloc"][bloc]["integrated"] += 1
        stats["by_bloc"][bloc]["items"].append(item)

        comp_str = item.get("COMPLÉTION_%", "0").rstrip("%").strip()
        try:
            comp = float(comp_str) if comp_str else 0
        except:
            comp = 0
        stats["by_bloc"][bloc]["completion_sum"] += comp

    # Calculate average per bloc
    for bloc in stats["by_bloc"]:
        total = stats["by_bloc"][bloc]["total"]
        if total > 0:
            stats["by_bloc"][bloc]["completion"] = stats["by_bloc"][bloc]["completion_sum"] / total
        else:
            stats["by_bloc"][bloc]["completion"] = 0

    return stats


def create_presentation():
    """Create base PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle, date_str):
    """Add title/cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["primary"]

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Date: {date_str}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER


def add_kpi_slide(prs, stats):
    """Add KPI/Executive Summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📊 EXECUTIVE SUMMARY"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # KPI Cards
    kpi_data = [
        ("INTÉGRÉ", f"{stats['integrated']}/{stats['total']}", COLORS["green"]),
        ("EN COURS", f"{stats['in_progress']}/{stats['total']}", COLORS["orange"]),
        ("PLANIFIÉ", f"{stats['planned']}/{stats['total']}", COLORS["secondary"]),
        ("BLOQUÉ", f"{stats['blocked']}/{stats['total']}", COLORS["red"]),
    ]

    y_position = 1.5
    for label, value, color in kpi_data:
        # Card background
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(y_position), Inches(9), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color

        # Label
        label_box = slide.shapes.add_textbox(Inches(1), Inches(y_position + 0.1), Inches(4), Inches(0.4))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

        # Value
        value_box = slide.shapes.add_textbox(Inches(5), Inches(y_position + 0.1), Inches(4), Inches(0.4))
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.RIGHT

        y_position += 1.2

    # Average completion
    avg_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    avg_frame = avg_box.text_frame
    p = avg_frame.paragraphs[0]
    p.text = f"⏱️  Complétion Moyenne: {stats['average_completion']:.1f}%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]


def add_completion_chart_slide(prs, stats):
    """Add completion by bloc bar chart slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📈 COMPLÉTION PAR BLOC"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # Add bars for each bloc
    y_position = 1.2
    for bloc, data in sorted(stats["by_bloc"].items(), key=lambda x: x[1]["completion"], reverse=True):
        completion = data["completion"]

        # Bloc label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_position), Inches(2.5), Inches(0.35))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = bloc[:25]  # Truncate long names
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]

        # Progress bar background
        bar_bg = slide.shapes.add_shape(1, Inches(3.2), Inches(y_position + 0.05), Inches(6), Inches(0.25))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(200, 200, 200)
        bar_bg.line.color.rgb = RGBColor(150, 150, 150)

        # Progress bar fill
        bar_width = (completion / 100) * 6
        bar = slide.shapes.add_shape(1, Inches(3.2), Inches(y_position + 0.05), Inches(bar_width), Inches(0.25))
        bar.fill.solid()
        if completion >= 90:
            bar.fill.fore_color.rgb = COLORS["green"]
        elif completion >= 70:
            bar.fill.fore_color.rgb = COLORS["accent"]
        else:
            bar.fill.fore_color.rgb = COLORS["orange"]
        bar.line.color.rgb = bar.fill.fore_color.rgb

        # Percentage
        pct_box = slide.shapes.add_textbox(Inches(9.3), Inches(y_position), Inches(0.5), Inches(0.35))
        pct_frame = pct_box.text_frame
        p = pct_frame.paragraphs[0]
        p.text = f"{completion:.0f}%"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.RIGHT

        y_position += 0.45


def add_gauge_slide(prs, stats):
    """Add overall completion gauge slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "🎯 COMPLÉTION GLOBALE DU PROJET"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # Large percentage display
    pct_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    pct_frame = pct_box.text_frame
    p = pct_frame.paragraphs[0]
    p.text = f"{stats['average_completion']:.0f}%"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER

    # Progress bar
    bar_bg = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(7), Inches(0.5))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(200, 200, 200)
    bar_bg.line.color.rgb = RGBColor(150, 150, 150)

    # Progress bar fill
    bar_width = (stats["average_completion"] / 100) * 7
    bar = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(bar_width), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["green"]
    bar.line.color.rgb = COLORS["green"]

    # Status text
    status_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    status_frame = status_box.text_frame
    status_frame.word_wrap = True
    p = status_frame.paragraphs[0]
    completion = stats["average_completion"]
    if completion >= 90:
        status = "🟢 Excellent - Quasi-finalisé"
    elif completion >= 70:
        status = "🟡 Bon - Sur la bonne voie"
    else:
        status = "🔴 À améliorer - Effort nécessaire"
    p.text = status
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER


def add_risks_slide(prs, risks):
    """Add risks and blockers slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "⚠️  RISQUES & BLOCAGES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    if not risks:
        no_risk_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        no_risk_frame = no_risk_box.text_frame
        p = no_risk_frame.paragraphs[0]
        p.text = "✅ Aucun risque ou blocage identifié"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS["green"]
        p.alignment = PP_ALIGN.CENTER
    else:
        y_position = 1.2
        for risk in risks[:5]:
            risk_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_position), Inches(8.4), Inches(0.9))
            risk_frame = risk_box.text_frame
            risk_frame.word_wrap = True

            element = risk.get("ÉLÉMENT", "Unknown")
            p = risk_frame.paragraphs[0]
            p.text = f"🟡 {element[:40]}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS["orange"]

            p = risk_frame.add_paragraph()
            statut = risk.get("STATUT", "")
            completion = risk.get("COMPLÉTION_%", "0")
            p.text = f"  Status: {statut} | Complétion: {completion}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["light_text"]

            y_position += 1.1


def add_next_steps_slide(prs):
    """Add next steps/recommendations slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📋 PROCHAINES ÉTAPES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    steps = [
        "✅ Phase 9: Mobile Responsive Design (3-5 jours)",
        "✅ Phase 10: E2E Tests avec Cypress (5-7 jours)",
        "✅ Phase 11: Configuration Email & Webhooks (2-3 jours)",
        "✅ Phase 12: Déploiement & Go-Live (2-3 jours)",
    ]

    y_position = 1.5
    for step in steps:
        step_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(8), Inches(0.8))
        step_frame = step_box.text_frame
        step_frame.word_wrap = True
        p = step_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(12)

        y_position += 0.9


def main():
    """Main execution"""
    print("\n🚀 PF Scoring - COPIL Presentation Generator (CSV Edition)")
    print("=" * 60)

    # Load and analyze data
    items = load_csv_data()
    stats = calculate_statistics(items)

    print(f"✅ Statistics calculated:")
    print(f"   - Total items: {stats['total']}")
    print(f"   - Integrated: {stats['integrated']}")
    print(f"   - In progress: {stats['in_progress']}")
    print(f"   - Average completion: {stats['average_completion']:.1f}%")

    # Create presentation
    print("\n📝 Creating presentation...")
    prs = create_presentation()

    # Add slides
    add_title_slide(prs, "PF SCORING V7++", "Rapport de Suivi - COPIL", datetime.now().strftime("%Y-%m-%d"))
    add_kpi_slide(prs, stats)
    add_completion_chart_slide(prs, stats)
    add_gauge_slide(prs, stats)
    add_risks_slide(prs, stats["risks"])
    add_next_steps_slide(prs)

    # Save presentation
    output_path = "COPIL_Presentation.pptx"
    prs.save(output_path)

    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n" + "=" * 60)


if __name__ == "__main__":
    main()
