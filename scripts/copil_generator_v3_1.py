#!/usr/bin/env python3
"""COPIL Banking v3.1.0 - Ultra Pro Edition with 9 Slides"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import csv
from datetime import datetime
import sys

# COLORS - Professional Banking Palette
PRIMARY_NAVY = RGBColor(0, 51, 102)
ACCENT_GOLD = RGBColor(255, 153, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(45, 45, 45)
LIGHT_GRAY = RGBColor(240, 240, 240)
LIGHT_BLUE = RGBColor(230, 240, 250)
VERY_LIGHT_BLUE = RGBColor(245, 250, 255)

GREEN = RGBColor(34, 139, 34)
LIGHT_GREEN = RGBColor(144, 238, 144)
YELLOW = RGBColor(255, 192, 0)
ORANGE = RGBColor(255, 128, 0)
RED = RGBColor(220, 20, 60)
LIGHT_RED = RGBColor(255, 200, 200)

def parse_completion(item):
    try:
        comp_str = item.get('COMPLÉTION_%', '0').rstrip('%').strip()
        return float(comp_str) if comp_str else 0
    except:
        return 0

def load_csv(csv_file):
    items = []
    with open(csv_file, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        items = list(reader)
    return items

def calc_stats(items):
    completions = [parse_completion(item) for item in items]
    risk_cats = {
        'Financial': len([i for i in items if 'financier' in i.get('CATÉGORIE', '').lower()]),
        'Technical': len([i for i in items if 'technique' in i.get('CATÉGORIE', '').lower()]),
        'Market': len([i for i in items if 'marché' in i.get('CATÉGORIE', '').lower()]),
        'Operational': len([i for i in items if 'opérationnel' in i.get('CATÉGORIE', '').lower()]),
    }
    stats = {
        'total': len(items),
        'integrated': len([i for i in items if i.get('STATUT') == 'INTÉGRÉ']),
        'in_progress': len([i for i in items if i.get('STATUT') == 'EN COURS']),
        'planned': len([i for i in items if i.get('STATUT') == 'PLANIFIÉ']),
        'blocked': len([i for i in items if i.get('STATUT') == 'BLOQUÉ']),
        'by_bloc': {},
        'risks': [i for i in items if i.get('STATUT') == 'BLOQUÉ' or parse_completion(i) < 50],
        'critical': [i for i in items if parse_completion(i) < 30],
        'risk_cats': risk_cats,
        'average_completion': sum(completions) / len(completions) if completions else 0,
    }
    for item in items:
        bloc = item.get('BLOC', 'Unknown')
        if bloc not in stats['by_bloc']:
            stats['by_bloc'][bloc] = {'total': 0, 'completion_sum': 0}
        stats['by_bloc'][bloc]['total'] += 1
        stats['by_bloc'][bloc]['completion_sum'] += parse_completion(item)
    for bloc in stats['by_bloc']:
        cnt = stats['by_bloc'][bloc]['total']
        stats['by_bloc'][bloc]['completion'] = stats['by_bloc'][bloc]['completion_sum'] / cnt if cnt > 0 else 0
    return stats

def add_title_bar(slide, text, icon=""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_NAVY
    bar.line.color.rgb = PRIMARY_NAVY

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.85), Inches(10), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{icon} {text}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

def create_metric_card(slide, x, y, value, label, color, sub_label=""):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(1.5), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.15), Inches(1.4), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.75), Inches(1.4), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if sub_label:
        tb = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 1.0), Inches(1.4), Inches(0.15))
        p = tb.text_frame.paragraphs[0]
        p.text = sub_label
        p.font.size = Pt(7)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def create_ppt(stats, output_file):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: COVER - Enhanced
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = PRIMARY_NAVY

    logo = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    logo.fill.solid()
    logo.fill.fore_color.rgb = ACCENT_GOLD

    logo_text = sl.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    p = logo_text.text_frame.paragraphs[0]
    p.text = "PF SCORING SYSTEM"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "PF SCORING"
    p.font.size = Pt(88)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Comité de Pilotage - COPIL"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = f"v3.1.0 | {datetime.now().strftime('%d/%m/%Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Conforme IFC • EBRD • Basel • Bank Al-Maghrib"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

    # SLIDE 2: EXECUTIVE SUMMARY - Enhanced
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "RÉSUMÉ EXÉCUTIF", "📋")

    comp = stats['average_completion']
    create_metric_card(sl, 1.2, 1.3, f"{comp:.0f}%", "Complétude", GREEN if comp >= 80 else ORANGE if comp >= 50 else RED, f"vs 100%")
    create_metric_card(sl, 2.95, 1.3, stats['integrated'], "Intégrés", GREEN, f"{stats['integrated']*100//stats['total']}%")
    create_metric_card(sl, 4.7, 1.3, stats['in_progress'], "En Cours", YELLOW, "Dev")
    create_metric_card(sl, 6.45, 1.3, stats['blocked'], "Bloqués", RED, "⚠️")

    tb = sl.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(8.6), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "STATUT GÉNÉRAL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY
    p.space_after = Pt(6)

    metrics = [
        f"📊 Projet: {stats['total']} éléments à livrer",
        f"✅ Livrés: {stats['integrated']} éléments ({stats['integrated']*100//stats['total']}%)",
        f"⚠️  En retard: {len(stats['risks'])} items ({len(stats['risks'])*100//stats['total']}%)",
        f"🔴 Critiques: {len(stats['critical'])} items < 30%",
        f"🎯 Objectif de complétude: 100% (Actuel: {comp:.0f}%)",
    ]

    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)

    # SLIDE 3: RISK ANALYSIS - Enhanced
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "ANALYSE DE RISQUE - IFC/EBRD", "⚠️")

    y = 1.3
    for cat, count in stats['risk_cats'].items():
        color = RED if count > 10 else ORANGE if count > 5 else GREEN

        label = sl.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.0), Inches(0.4))
        p = label.text_frame.paragraphs[0]
        p.text = f"• {cat}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY

        count_box = sl.shapes.add_textbox(Inches(2.8), Inches(y), Inches(0.6), Inches(0.4))
        p = count_box.text_frame.paragraphs[0]
        p.text = str(count)
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        bar_bg = sl.shapes.add_shape(1, Inches(3.6), Inches(y + 0.08), Inches(5.7), Inches(0.25))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_GRAY

        bar_fill = sl.shapes.add_shape(1, Inches(3.6), Inches(y + 0.08), Inches(5.7 * min(count / 15, 1)), Inches(0.25))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color

        y += 1.35

    rec = sl.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.6), Inches(0.9))
    rec.text_frame.word_wrap = True
    p = rec.text_frame.paragraphs[0]
    p.text = "🔴 Action Requise: Escalade pour éléments bloqués et conformité réglementaire immédiate"
    p.font.size = Pt(11)
    p.font.color.rgb = RED
    p.font.bold = True

    # SLIDE 4: KPI DASHBOARD - Enhanced
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "TABLEAU DE BORD - KPIs", "📊")

    kpis = [
        ("INTÉGRÉ", stats['integrated'], GREEN, "Livré"),
        ("EN COURS", stats['in_progress'], YELLOW, "Dev"),
        ("PLANIFIÉ", stats['planned'], ORANGE, "Prévu"),
        ("BLOQUÉ", stats['blocked'], RED, "🚫"),
    ]

    for idx, (lbl, val, col, desc) in enumerate(kpis):
        x = 0.6 + (idx * 2.3)
        card = sl.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(2), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = col
        card.line.width = Pt(2)

        tb = sl.shapes.add_textbox(Inches(x + 0.1), Inches(1.5), Inches(1.8), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = str(val)
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = sl.shapes.add_textbox(Inches(x + 0.1), Inches(2.5), Inches(1.8), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = sl.shapes.add_textbox(Inches(x + 0.1), Inches(2.85), Inches(1.8), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(8.6), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "INDICATEURS CLÉS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY
    p.space_after = Pt(4)

    kpi_items = [
        f"Taux de complétion: {stats['average_completion']:.1f}% (Target: 100%)",
        f"Velocity: {stats['integrated']} livrés / {stats['total']} total ({stats['integrated']*100//stats['total']}%)",
        f"Items en retard: {len(stats['risks'])} ({len(stats['risks'])*100//stats['total']}% du projet)",
        f"Items critiques: {len(stats['critical'])} éléments < 30% complet",
        f"Blockers identifiés: {stats['blocked']} items nécessitant escalade",
    ]

    for item in kpi_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)

    # SLIDE 5: ACTION PLAN - Detailed
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "PLAN D'ACTION STRATÉGIQUE", "✅")

    actions = [
        ("🔴 IMMÉDIAT (0-2 sem)", "• Escalade items bloqués\n• Review risques critiques\n• Validation governance", RED),
        ("🟠 COURT TERME (2-4 sem)", "• Accélération items en retard\n• Allocation ressources\n• Replan si nécessaire", ORANGE),
        ("🟡 MOYEN TERME (1-2 mois)", "• Completion conformité\n• Testing & QA\n• Préparation déploiement", YELLOW),
        ("🟢 LONG TERME (2+ mois)", "• Déploiement production\n• Monitoring post-launch\n• Support utilisateurs", GREEN),
    ]

    y = 1.3
    for title, actions_text, color in actions:
        box = sl.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(8.8), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = color
        box.line.width = Pt(3)

        tb = sl.shapes.add_textbox(Inches(0.8), Inches(y + 0.05), Inches(8.4), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        tb = sl.shapes.add_textbox(Inches(0.8), Inches(y + 0.4), Inches(8.4), Inches(0.85))
        p = tb.text_frame.paragraphs[0]
        p.text = actions_text
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GRAY

        y += 1.5

    # SLIDE 6: COMPLETION GAUGE
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "AVANCEMENT GLOBAL", "📈")

    comp = stats['average_completion']
    tb = sl.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{comp:.0f}%"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = GREEN if comp >= 80 else ORANGE if comp >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    status = "✅ EXCELLENT" if comp >= 80 else "⚠️  ACCEPTABLE" if comp >= 50 else "🔴 À AMÉLIORER"
    tb = sl.shapes.add_textbox(Inches(2), Inches(3.1), Inches(6), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = status
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN if comp >= 80 else ORANGE if comp >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    bar_bg = sl.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(7), Inches(0.5))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = LIGHT_GRAY

    bar_fg = sl.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(7 * (comp / 100)), Inches(0.5))
    bar_fg.fill.solid()
    bar_fg.fill.fore_color.rgb = GREEN if comp >= 80 else ORANGE if comp >= 50 else RED

    tb = sl.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True

    metrics = [
        f"Total: {stats['total']} éléments",
        f"Intégrés: {stats['integrated']} ({stats['integrated']*100//stats['total']}%)",
        f"En retard: {len(stats['risks'])} items",
        f"Critiques: {len(stats['critical'])} items < 30%",
    ]

    for i, metric in enumerate(metrics):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)

    # SLIDE 7: DÉTAILS PAR BLOC - NEW
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "DÉTAILS PAR BLOC", "📦")

    tb = sl.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "BLOC"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY
    p.space_after = Pt(2)

    for bloc, data in list(stats['by_bloc'].items())[:4]:
        p = tf.add_paragraph()
        p.text = f"  • {bloc}: {data['completion']:.0f}% ({data['total']} items)"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)

    # SLIDE 8: RECOMMENDATIONS - NEW
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    add_title_bar(sl, "RECOMMANDATIONS", "💡")

    tb = sl.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True

    recommendations = [
        "1️⃣ PRIORITÉ 1: Résoudre les {0} items bloqués - Impact: Critique - Timeline: IMMÉDIAT",
        "2️⃣ PRIORITÉ 2: Accélérer les {0} items en retard - Impact: Haut - Timeline: 2-4 semaines",
        "3️⃣ PRIORITÉ 3: Identifier et mitiguer les {0} risques critiques - Impact: Moyen - Timeline: En cours",
        "🎯 OBJECTIF: Atteindre 100% de complétude avant déploiement production",
        "📅 CHECKPOINT: Revoir statut chaque semaine en COPIL",
    ]

    rec_text = "\n".join(recommendations).format(stats['blocked'], len(stats['risks']), len(stats['critical']))

    p = tf.paragraphs[0]
    p.text = rec_text
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

    # SLIDE 9: CLOSING PRO
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = PRIMARY_NAVY

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Merci"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = f"PF Scoring v3.1.0 | {datetime.now().strftime('%d/%m/%Y')} | IFC • EBRD • Basel • BAM"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_file)
    print(f'✅ SUCCESS: {len(prs.slides)} slides ULTRA PRO générées! (v3.1.0)')

if __name__ == '__main__':
    csv_file = sys.argv[1] if len(sys.argv) > 1 else 'PF_SCORING_SPECIFICATIONS_TRACKING.csv'
    output_file = sys.argv[2] if len(sys.argv) > 2 else 'COPIL_Banking_Presentation_v3_1.pptx'

    try:
        items = load_csv(csv_file)
        print(f'✅ Chargés: {len(items)} éléments')
        stats = calc_stats(items)
        create_ppt(stats, output_file)
        print(f'📁 Fichier généré: {output_file}')
    except FileNotFoundError:
        print(f'❌ Fichier CSV introuvable: {csv_file}')
        sys.exit(1)
    except Exception as e:
        print(f'❌ Erreur: {e}')
        sys.exit(1)
