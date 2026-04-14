#!/usr/bin/env python3
"""
COPIL Presentation Generator
Version: 2.0.0

Generates professional PowerPoint presentations for COPIL (Governance Committee) meetings
from PF Scoring tracking CSV data.

Features:
  - 6-slide professional presentations
  - KPI summaries (INTÉGRÉ/EN COURS/PLANIFIÉ/BLOQUÉ)
  - Completion charts by bloc with status coloring
  - Global completion gauge
  - Risk and blocker identification
  - Professional banking color scheme

Usage:
  python copil_generator.py <csv_file> [output_file]

Examples:
  python copil_generator.py PF_SCORING_SPECIFICATIONS_TRACKING.csv
  python copil_generator.py data.csv COPIL_Report.pptx

Author: Claude Code
License: Internal Use Only
"""

__version__ = "2.1.0"
__author__ = "Claude Code"
__date__ = "2026-04-06"

import csv
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE


# ═══════════════════════════════════════════════════════════════════════════
# COLOR SCHEME - Professional Banking Theme
# ═══════════════════════════════════════════════════════════════════════════
PRIMARY_BLUE = RGBColor(0, 51, 102)        # #003366 - Dark blue
ACCENT_ORANGE = RGBColor(255, 153, 0)     # #FF9900 - Accent orange
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(76, 175, 80)             # Success/Good (≥90%)
ORANGE = RGBColor(255, 152, 0)            # Warning/Acceptable (≥70%)
RED = RGBColor(244, 67, 54)               # Alert/Poor (<70%)
DARK_TEXT = RGBColor(33, 33, 33)
LIGHT_GRAY = RGBColor(240, 240, 240)


# ═══════════════════════════════════════════════════════════════════════════
# UTILITIES
# ═══════════════════════════════════════════════════════════════════════════

def parse_completion(item: dict) -> float:
    """
    Safely parse completion percentage from item dictionary.

    Args:
        item: Dictionary with 'COMPLÉTION_%' key

    Returns:
        Float between 0-100, defaults to 0 if parsing fails
    """
    try:
        comp_str = item.get('COMPLÉTION_%', '0').rstrip('%').strip()
        return float(comp_str) if comp_str else 0
    except (ValueError, AttributeError):
        return 0


def load_csv_data(csv_file: str) -> list:
    """
    Load CSV data into list of dictionaries.

    Args:
        csv_file: Path to CSV file

    Returns:
        List of dictionaries with CSV data

    Raises:
        FileNotFoundError: If CSV file doesn't exist
        ValueError: If CSV is invalid
    """
    if not os.path.exists(csv_file):
        raise FileNotFoundError(f"CSV file not found: {csv_file}")

    print(f"📥 Loading: {os.path.basename(csv_file)}")
    items = []

    try:
        with open(csv_file, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            items = list(reader)
    except Exception as e:
        raise ValueError(f"Failed to parse CSV: {str(e)}")

    if not items:
        raise ValueError("CSV file is empty")

    print(f"✅ Loaded {len(items)} items")
    return items


def calculate_statistics(items: list) -> dict:
    """
    Calculate KPIs and statistics from items.

    Args:
        items: List of tracking items

    Returns:
        Dictionary with statistics:
          - total: Total number of items
          - integrated: Count of INTÉGRÉ items
          - in_progress: Count of EN COURS items
          - planned: Count of PLANIFIÉ items
          - blocked: Count of BLOQUÉ items
          - average_completion: Average completion percentage
          - by_bloc: Breakdown by bloc
          - risks: List of risky items (BLOQUÉ or <50% complete)
    """
    print("📊 Calculating statistics...")

    completions = [parse_completion(item) for item in items]

    stats = {
        'total': len(items),
        'integrated': len([i for i in items if i.get('STATUT') == 'INTÉGRÉ']),
        'in_progress': len([i for i in items if i.get('STATUT') == 'EN COURS']),
        'planned': len([i for i in items if i.get('STATUT') == 'PLANIFIÉ']),
        'blocked': len([i for i in items if i.get('STATUT') == 'BLOQUÉ']),
        'by_bloc': {},
        'risks': [i for i in items if i.get('STATUT') == 'BLOQUÉ' or parse_completion(i) < 50]
    }

    stats['average_completion'] = sum(completions) / len(completions) if completions else 0

    # Group by bloc
    for item in items:
        bloc = item.get('BLOC', 'Unknown')
        if bloc not in stats['by_bloc']:
            stats['by_bloc'][bloc] = {'total': 0, 'completion_sum': 0}
        stats['by_bloc'][bloc]['total'] += 1
        stats['by_bloc'][bloc]['completion_sum'] += parse_completion(item)

    # Calculate completion per bloc
    for bloc in stats['by_bloc']:
        count = stats['by_bloc'][bloc]['total']
        stats['by_bloc'][bloc]['completion'] = (
            stats['by_bloc'][bloc]['completion_sum'] / count if count > 0 else 0
        )

    print(f"✅ {stats['total']} items | {stats['average_completion']:.1f}% complete")
    return stats


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE CREATORS
# ═══════════════════════════════════════════════════════════════════════════

def add_title_slide(prs: Presentation, stats: dict) -> None:
    """Create professional title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY_BLUE

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_ORANGE
    accent_bar.line.color.rgb = ACCENT_ORANGE

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "PF SCORING"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.2))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "COMITÉ DE PILOTAGE"
    p.font.size = Pt(44)
    p.font.color.rgb = ACCENT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Sub-subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.95), Inches(9), Inches(0.6))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Présentation COPIL"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    line = slide.shapes.add_shape(1, Inches(3.5), Inches(5.8), Inches(3), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.color.rgb = ACCENT_ORANGE

    # Status badge
    badge = slide.shapes.add_shape(1, Inches(3), Inches(6.2), Inches(4), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_ORANGE
    badge.line.color.rgb = ACCENT_ORANGE

    status_box = slide.shapes.add_textbox(Inches(3), Inches(6.25), Inches(4), Inches(0.5))
    p = status_box.text_frame.paragraphs[0]
    p.text = f"Avancement: {stats['average_completion']:.0f}%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    p.alignment = PP_ALIGN.CENTER


def add_kpi_slide(prs: Presentation, stats: dict) -> None:
    """Create KPI summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_BLUE
    title_shape.line.color.rgb = PRIMARY_BLUE

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "RÉSUMÉ EXÉCUTIF - KPIs"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # KPI cards with better spacing
    kpis = [
        ("INTÉGRÉ", stats['integrated'], GREEN, f"{stats['integrated']} items"),
        ("EN COURS", stats['in_progress'], ORANGE, f"{stats['in_progress']} items"),
        ("PLANIFIÉ", stats['planned'], RGBColor(33, 150, 243), f"{stats['planned']} items"),
        ("BLOQUÉ", stats['blocked'], RED, f"{stats['blocked']} items")
    ]

    for idx, (label, count, color, desc) in enumerate(kpis):
        x = 0.7 + (idx * 2.2)
        y = 1.3

        # Card shape with shadow effect
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        # Count - Large
        count_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.4), Inches(1.8), Inches(0.8))
        p = count_box.text_frame.paragraphs[0]
        p.text = str(count)
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.3), Inches(1.8), Inches(0.35))
        p = label_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.7), Inches(1.8), Inches(0.2))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Summary bar
    summary_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    summary_shape.fill.solid()
    summary_shape.fill.fore_color.rgb = LIGHT_GRAY
    summary_shape.line.color.rgb = LIGHT_GRAY

    summary_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.65), Inches(8.6), Inches(1.2))
    p = summary_box.text_frame.paragraphs[0]
    p.text = f"Total des éléments: {stats['total']} | Complétude moyenne: {stats['average_completion']:.1f}% | Intégrés: {stats['integrated']}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    p.alignment = PP_ALIGN.CENTER
    summary_box.text_frame.word_wrap = True


def add_completion_chart_slide(prs: Presentation, stats: dict) -> None:
    """Create completion chart by bloc."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "AVANCEMENT PAR BLOC"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Chart
    try:
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.1), Inches(9), Inches(5.2)
        ).chart
        chart.has_legend = False

        categories = list(stats['by_bloc'].keys())
        completions = tuple([stats['by_bloc'][bloc]['completion'] for bloc in categories])

        chart.chart_data.categories = categories
        chart.chart_data.add_series('Complétude %', completions)

        # Color series based on completion
        series = chart.plots[0].series[0]
        for idx, completion in enumerate(completions):
            if completion >= 90:
                color = GREEN
            elif completion >= 70:
                color = ORANGE
            else:
                color = RED

            series.points[idx].format.fill.solid()
            series.points[idx].format.fill.fore_color.rgb = color
    except Exception as e:
        print(f"⚠️  Chart creation skipped: {str(e)}")


def add_gauge_slide(prs: Presentation, stats: dict) -> None:
    """Create global completion gauge with progress bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_BLUE
    title_shape.line.color.rgb = PRIMARY_BLUE

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "COMPLÉTUDE GLOBALE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    completion = stats['average_completion']

    # Main percentage display
    pct_box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1.5))
    p = pct_box.text_frame.paragraphs[0]
    p.text = f"{completion:.0f}%"
    p.font.size = Pt(140)
    p.font.bold = True
    p.font.color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    # Status text
    status_text = "✅ EXCELLENT" if completion >= 80 else "⚠️  ACCEPTABLE" if completion >= 50 else "❌ À AMÉLIORER"
    status_box = slide.shapes.add_textbox(Inches(2), Inches(3.4), Inches(6), Inches(0.4))
    p = status_box.text_frame.paragraphs[0]
    p.text = status_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    # Progress bar background
    bar_bg = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(7), Inches(0.5))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = LIGHT_GRAY
    bar_bg.line.color.rgb = RGBColor(180, 180, 180)

    # Progress bar fill
    bar_fill = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(7 * (completion / 100)), Inches(0.5))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    bar_fill.line.color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED

    # Stats below
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.8))
    stats_box.text_frame.word_wrap = True
    p = stats_box.text_frame.paragraphs[0]
    p.text = f"Éléments INTÉGRÉS: {stats['integrated']} | EN COURS: {stats['in_progress']} | PLANIFIÉS: {stats['planned']} | BLOQUÉS: {stats['blocked']}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    p2 = stats_box.text_frame.add_paragraph()
    p2.text = f"Total: {stats['total']} éléments"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER


def add_risks_slide(prs: Presentation, stats: dict) -> None:
    """Create risks and blockers slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "RISQUES & BLOCAGES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    risks = stats['risks'][:5]

    if not risks:
        no_risks = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        p = no_risks.text_frame.paragraphs[0]
        p.text = "✅ Aucun risque majeur identifié"
        p.font.size = Pt(28)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
    else:
        y = 1.3
        for risk in risks:
            box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.8))
            p = box.text_frame.paragraphs[0]
            element = risk.get('ÉLÉMENT', 'Unknown')
            status = risk.get('STATUT', '')
            comp = parse_completion(risk)
            p.text = f"🔴 {element} ({status} - {comp:.0f}%)"
            p.font.size = Pt(14)
            p.font.color.rgb = RED
            y += 0.9


def add_next_steps_slide(prs: Presentation, stats: dict) -> None:
    """Create next steps timeline slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = "PROCHAINES ÉTAPES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    phases = [
        ("Phase 9", "Design Responsive Mobile"),
        ("Phase 10", "Tests & QA Complète"),
        ("Phase 11", "Déploiement Production"),
        ("Phase 12", "Monitoring & Support"),
    ]

    y = 1.3
    for phase, description in phases:
        # Phase marker
        shape = slide.shapes.add_shape(1, Inches(0.7), Inches(y), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_ORANGE
        shape.line.color.rgb = ACCENT_ORANGE

        # Description
        box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.05), Inches(7.8), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = f"{phase} - {description}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT

        y += 1


def create_presentation(stats: dict, output_file: str) -> None:
    """
    Create and save complete presentation.

    Args:
        stats: Statistics dictionary from calculate_statistics()
        output_file: Path where to save the PPTX file
    """
    print("\n📝 Creating presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add all slides
    add_title_slide(prs, stats)
    add_kpi_slide(prs, stats)
    add_completion_chart_slide(prs, stats)
    add_gauge_slide(prs, stats)
    add_risks_slide(prs, stats)
    add_next_steps_slide(prs, stats)

    # Save
    prs.save(output_file)

    print(f"\n✅ Presentation created successfully!")
    print(f"📁 File: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    """Main entry point."""
    print("\n" + "="*60)
    print("🚀 PF Scoring - COPIL Presentation Generator")
    print(f"Version: {__version__}")
    print("="*60)

    # Parse arguments
    if len(sys.argv) < 2:
        print("\nUsage: python copil_generator.py <csv_file> [output_file]")
        print("\nExample:")
        print("  python copil_generator.py PF_SCORING_SPECIFICATIONS_TRACKING.csv")
        print("  python copil_generator.py data.csv COPIL_Report.pptx")
        sys.exit(1)

    csv_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "COPIL_Presentation.pptx"

    try:
        # Load and process
        items = load_csv_data(csv_file)
        stats = calculate_statistics(items)
        create_presentation(stats, output_file)

        print("="*60)
        print("✅ SUCCESS!\n")

    except FileNotFoundError as e:
        print(f"\n❌ Error: {str(e)}")
        sys.exit(1)
    except ValueError as e:
        print(f"\n❌ Error: {str(e)}")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ Unexpected error: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()
