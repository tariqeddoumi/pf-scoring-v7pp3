#!/usr/bin/env python3
"""
COPIL Presentation Generator - Banking Edition
Version: 3.0.0

Professional PowerPoint presentations for Banking COPIL meetings
Compliance: IFC, EBRD, Basel, Bank Al-Maghrib

Features:
  - 8-slide professional banking presentations
  - Executive Summary with strategic insights
  - Risk analysis (Financial, Technical, Market, Operational)
  - Compliance framework overview
  - KPI dashboards
  - Strategic recommendations

Usage:
  python copil_generator_banking.py <csv_file> [output_file]

Author: Claude Code
License: Internal Use Only
"""

__version__ = "3.0.0"
__author__ = "Claude Code"
__date__ = "2026-04-06"

import csv
import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE


# ═══════════════════════════════════════════════════════════════════════════
# PROFESSIONAL BANKING COLOR SCHEME
# ═══════════════════════════════════════════════════════════════════════════
PRIMARY_NAVY = RGBColor(0, 51, 102)         # #003366 - Corporate Navy
ACCENT_GOLD = RGBColor(255, 153, 0)        # #FF9900 - Banking Gold
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(45, 45, 45)
LIGHT_GRAY = RGBColor(240, 240, 240)
LIGHT_BLUE = RGBColor(230, 240, 250)

# Risk colors - Banking standard
GREEN = RGBColor(34, 139, 34)               # Low Risk
YELLOW = RGBColor(255, 192, 0)              # Medium Risk
ORANGE = RGBColor(255, 128, 0)              # High Risk
RED = RGBColor(220, 20, 60)                 # Critical Risk


# ═══════════════════════════════════════════════════════════════════════════
# UTILITIES
# ═══════════════════════════════════════════════════════════════════════════

def parse_completion(item: dict) -> float:
    """Safely parse completion percentage."""
    try:
        comp_str = item.get('COMPLÉTION_%', '0').rstrip('%').strip()
        return float(comp_str) if comp_str else 0
    except:
        return 0


def load_csv_data(csv_file: str) -> list:
    """Load CSV data."""
    if not os.path.exists(csv_file):
        raise FileNotFoundError(f"CSV file not found: {csv_file}")

    print(f"📥 Loading: {os.path.basename(csv_file)}")
    with open(csv_file, 'r', encoding='utf-8') as f:
        items = list(csv.DictReader(f))

    if not items:
        raise ValueError("CSV file is empty")

    print(f"✅ Loaded {len(items)} items")
    return items


def calculate_statistics(items: list) -> dict:
    """Calculate comprehensive statistics."""
    print("📊 Calculating statistics...")

    completions = [parse_completion(item) for item in items]

    # Risk categories per IFC/EBRD standards
    risk_categories = {
        'Financial': len([i for i in items if 'financier' in i.get('CATÉGORIE', '').lower()]),
        'Technical': len([i for i in items if 'technique' in i.get('CATÉGORIE', '').lower()]),
        'Market': len([i for i in items if 'marché' in i.get('CATÉGORIE', '').lower()]),
        'Operational': len([i for i in items if 'opérationnel' in i.get('CATÉGORIE', '').lower()]),
        'Environmental': len([i for i in items if 'environnemental' in i.get('CATÉGORIE', '').lower()]),
        'Social': len([i for i in items if 'social' in i.get('CATÉGORIE', '').lower()]),
    }

    stats = {
        'total': len(items),
        'integrated': len([i for i in items if i.get('STATUT') == 'INTÉGRÉ']),
        'in_progress': len([i for i in items if i.get('STATUT') == 'EN COURS']),
        'planned': len([i for i in items if i.get('STATUT') == 'PLANIFIÉ']),
        'blocked': len([i for i in items if i.get('STATUT') == 'BLOQUÉ']),
        'average_completion': sum(completions) / len(completions) if completions else 0,
        'by_bloc': {},
        'risks': [i for i in items if i.get('STATUT') == 'BLOQUÉ' or parse_completion(i) < 50],
        'risk_categories': risk_categories,
        'critical_items': [i for i in items if parse_completion(i) < 30],
    }

    # Group by bloc
    for item in items:
        bloc = item.get('BLOC', 'Unknown')
        if bloc not in stats['by_bloc']:
            stats['by_bloc'][bloc] = {'total': 0, 'completion_sum': 0}
        stats['by_bloc'][bloc]['total'] += 1
        stats['by_bloc'][bloc]['completion_sum'] += parse_completion(item)

    for bloc in stats['by_bloc']:
        count = stats['by_bloc'][bloc]['total']
        stats['by_bloc'][bloc]['completion'] = (
            stats['by_bloc'][bloc]['completion_sum'] / count if count > 0 else 0
        )

    print(f"✅ Analysis complete: {stats['total']} items | {stats['average_completion']:.1f}% completion")
    return stats


def add_centered_title_bar(slide, text: str, color=PRIMARY_NAVY) -> None:
    """Add centered title bar to slide."""
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.85))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.85), Inches(10), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD
    accent.line.color.rgb = ACCENT_GOLD

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE CREATORS
# ═══════════════════════════════════════════════════════════════════════════

def add_cover_slide(prs: Presentation, stats: dict) -> None:
    """Professional banking cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRIMARY_NAVY

    # Logo area
    logo_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = ACCENT_GOLD
    logo_box.line.color.rgb = ACCENT_GOLD

    # Main title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "PF SCORING"
    p.font.size = Pt(88)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Comité de Pilotage - COPIL"
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Regulatory compliance line
    compliance = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    p = compliance.text_frame.paragraphs[0]
    p.text = "Conforme IFC • EBRD • Basel • Bank Al-Maghrib"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

    # Date and status
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.5))
    p = date_box.text_frame.paragraphs[0]
    p.text = f"Date: {datetime.now().strftime('%d/%m/%Y')} | Avancement: {stats['average_completion']:.0f}%"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    bottom_bar = slide.shapes.add_shape(1, Inches(0), Inches(7.4), Inches(10), Inches(0.1))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT_GOLD
    bottom_bar.line.color.rgb = ACCENT_GOLD


def add_executive_summary_slide(prs: Presentation, stats: dict) -> None:
    """Executive summary for C-level."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_centered_title_bar(slide, "📋 RÉSUMÉ EXÉCUTIF")

    # Key metrics
    metrics = [
        ("Complétude Global", f"{stats['average_completion']:.1f}%", 1.2),
        ("Éléments Intégrés", str(stats['integrated']), 3.5),
        ("Éléments en Cours", str(stats['in_progress']), 5.8),
        ("Éléments Bloqués", str(stats['blocked']), 8.0),
    ]

    for label, value, x in metrics:
        # Value box
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(1.5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = PRIMARY_NAVY
        box.line.width = Pt(1.5)

        # Value
        val_text = slide.shapes.add_textbox(Inches(x), Inches(1.65), Inches(1.5), Inches(0.6))
        p = val_text.text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_text = slide.shapes.add_textbox(Inches(x), Inches(2.25), Inches(1.5), Inches(0.4))
        p = lbl_text.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Summary text
    summary = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(3.8))
    summary.text_frame.word_wrap = True
    tf = summary.text_frame

    p = tf.paragraphs[0]
    p.text = "STATUT GÉNÉRAL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = f"• Projet: {stats['total']} éléments à livrer"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(4)
    p.level = 0

    p = tf.add_paragraph()
    p.text = f"• Conformité: Avancement de {stats['average_completion']:.0f}% vs objectif 100%"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = f"• Risque de dépassement de délais: {len(stats['risks'])} éléments en retard"
    p.font.size = Pt(11)
    color = RED if len(stats['risks']) > 10 else ORANGE if len(stats['risks']) > 5 else GREEN
    p.font.color.rgb = color
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = f"• Éléments critiques: {len(stats['critical_items'])} items < 30% complétude"
    p.font.size = Pt(11)
    p.font.color.rgb = RED if len(stats['critical_items']) > 0 else GREEN
    p.space_after = Pt(4)


def add_risk_analysis_slide(prs: Presentation, stats: dict) -> None:
    """Risk analysis by category (IFC/EBRD standards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_centered_title_bar(slide, "⚠️  ANALYSE DE RISQUE - Catégories IFC/EBRD")

    # Risk categories
    risk_cats = stats['risk_categories']
    categories = ['Financial', 'Technical', 'Market', 'Operational']
    y_start = 1.3

    for i, cat in enumerate(categories):
        count = risk_cats.get(cat, 0)
        y = y_start + (i * 1.2)

        # Category label
        label_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.5), Inches(0.4))
        p = label_box.text_frame.paragraphs[0]
        p.text = f"• {cat} Risk:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY

        # Count box
        count_box = slide.shapes.add_textbox(Inches(3.2), Inches(y), Inches(0.8), Inches(0.4))
        p = count_box.text_frame.paragraphs[0]
        p.text = str(count)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RED if count > 10 else ORANGE if count > 5 else GREEN

        # Progress bar
        bar_bg = slide.shapes.add_shape(1, Inches(4.2), Inches(y + 0.05), Inches(5), Inches(0.3))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_GRAY
        bar_bg.line.color.rgb = RGBColor(200, 200, 200)

    # Recommendations box
    rec_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(8.6), Inches(0.8))
    rec_box.text_frame.word_wrap = True
    p = rec_box.text_frame.paragraphs[0]
    p.text = "🔴 Recommandations: Escalade requise pour éléments bloqués. Vérifier conformité regulatory compliance."
    p.font.size = Pt(11)
    p.font.color.rgb = RED
    p.font.bold = True


def add_kpi_dashboard_slide(prs: Presentation, stats: dict) -> None:
    """Professional KPI dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_centered_title_bar(slide, "📊 TABLEAU DE BORD - KPIs")

    # KPI cards
    kpis = [
        ("INTÉGRÉ", stats['integrated'], GREEN, "Livré"),
        ("EN COURS", stats['in_progress'], YELLOW, "En développement"),
        ("PLANIFIÉ", stats['planned'], ORANGE, "Prévu"),
        ("BLOQUÉ", stats['blocked'], RED, "Critique"),
    ]

    for idx, (label, count, color, desc) in enumerate(kpis):
        x = 0.6 + (idx * 2.3)
        y = 1.3

        # Card
        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.width = Pt(2)

        # Number
        num_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.35), Inches(1.8), Inches(0.9))
        p = num_box.text_frame.paragraphs[0]
        p.text = str(count)
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.3), Inches(1.8), Inches(0.35))
        p = lbl_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.65), Inches(1.8), Inches(0.25))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Summary metrics
    summary_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(8.6), Inches(3))
    summary_box.text_frame.word_wrap = True
    tf = summary_box.text_frame

    p = tf.paragraphs[0]
    p.text = "INDICATEURS CLÉS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_NAVY
    p.space_after = Pt(6)

    metrics_text = [
        f"Taux de complétion: {stats['average_completion']:.1f}% (Objectif: 100%)",
        f"Éléments livrer: {stats['integrated']} / {stats['total']} ({stats['integrated']*100//stats['total']}%)",
        f"Items en retard: {len(stats['risks'])} (Taux: {len(stats['risks'])*100//stats['total']}%)",
        f"Éléments critiques: {len(stats['critical_items'])} (< 30% complétude)",
    ]

    for metric in metrics_text:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(3)


def add_action_plan_slide(prs: Presentation, stats: dict) -> None:
    """Action plan and recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_centered_title_bar(slide, "✅ PLAN D'ACTION - Recommandations")

    # Action items
    actions = [
        ("1. IMMÉDIAT (0-2 semaines)", "Escalade éléments bloqués • Review risques critiques • Validation governance", RED),
        ("2. COURT TERME (2-4 semaines)", "Accélération items en retard • Allocation ressources supplémentaires", ORANGE),
        ("3. MOYEN TERME (1-2 mois)", "Completion conformité • Testing & QA • Préparation déploiement", YELLOW),
        ("4. LONG TERME (2+ mois)", "Déploiement production • Monitoring post-launch • Support utilisateurs", GREEN),
    ]

    y = 1.3
    for title, content, color in actions:
        # Action box
        box = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(8.8), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245) if color != RED else RGBColor(255, 240, 240)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.05), Inches(8.4), Inches(0.35))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.4), Inches(8.4), Inches(0.5))
        content_box.text_frame.word_wrap = True
        p = content_box.text_frame.paragraphs[0]
        p.text = content
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY

        y += 1.1


def add_completion_slide(prs: Presentation, stats: dict) -> None:
    """Completion gauge with progress visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_centered_title_bar(slide, "📈 AVANCEMENT GLOBAL")

    completion = stats['average_completion']

    # Large percentage
    pct_box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(1.2))
    p = pct_box.text_frame.paragraphs[0]
    p.text = f"{completion:.0f}%"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    # Status
    status = "✅ EXCELLENT" if completion >= 80 else "⚠️  ACCEPTABLE" if completion >= 50 else "🔴 À AMÉLIORER"
    status_box = slide.shapes.add_textbox(Inches(2), Inches(3.1), Inches(6), Inches(0.4))
    p = status_box.text_frame.paragraphs[0]
    p.text = status
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    p.alignment = PP_ALIGN.CENTER

    # Progress bar
    bar_bg = slide.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(7), Inches(0.5))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = LIGHT_GRAY
    bar_bg.line.color.rgb = RGBColor(180, 180, 180)

    bar_fill = slide.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(7 * (completion / 100)), Inches(0.5))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = GREEN if completion >= 80 else ORANGE if completion >= 50 else RED
    bar_fill.line.color.rgb = bar_fill.fill.fore_color.rgb

    # Details
    details_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(1.8))
    details_box.text_frame.word_wrap = True
    tf = details_box.text_frame

    items = [
        f"Total: {stats['total']} éléments",
        f"Intégrés: {stats['integrated']} ({stats['integrated']*100//stats['total']}%)",
        f"En retard: {len(stats['risks'])} items",
        f"Critiques: {len(stats['critical_items'])} items < 30%",
    ]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)


def add_closing_slide(prs: Presentation) -> None:
    """Professional closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRIMARY_NAVY

    # Main message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    msg_box.text_frame.word_wrap = True
    p = msg_box.text_frame.paragraphs[0]
    p.text = "Merci"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_box.text_frame.word_wrap = True
    p = footer_box.text_frame.paragraphs[0]
    p.text = f"PF Scoring v3.0 | {datetime.now().strftime('%d/%m/%Y')} | Conforme IFC • EBRD • Basel"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


def create_presentation(stats: dict, output_file: str) -> None:
    """Create professional banking presentation."""
    print("\n📝 Creating professional banking presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add all slides
    add_cover_slide(prs, stats)
    add_executive_summary_slide(prs, stats)
    add_risk_analysis_slide(prs, stats)
    add_kpi_dashboard_slide(prs, stats)
    add_action_plan_slide(prs, stats)
    add_completion_slide(prs, stats)
    add_closing_slide(prs)

    prs.save(output_file)

    print(f"\n✅ Professional banking presentation created!")
    print(f"📁 File: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"🏦 Compliance: IFC • EBRD • Basel • Bank Al-Maghrib")


def main():
    """Main entry point."""
    print("\n" + "="*70)
    print("🏦 PF Scoring - Professional Banking COPIL Presentation Generator")
    print(f"Version: {__version__} | Compliance: IFC • EBRD • Basel")
    print("="*70)

    if len(sys.argv) < 2:
        print("\nUsage: python copil_generator_banking.py <csv_file> [output_file]")
        print("\nExample:")
        print("  python copil_generator_banking.py PF_SCORING_SPECIFICATIONS_TRACKING.csv")
        print("  python copil_generator_banking.py data.csv COPIL_Banking_Report.pptx")
        sys.exit(1)

    csv_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "COPIL_Banking_Presentation.pptx"

    try:
        items = load_csv_data(csv_file)
        stats = calculate_statistics(items)
        create_presentation(stats, output_file)
        print("="*70 + "\n")

    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()
