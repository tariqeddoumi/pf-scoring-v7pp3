#!/usr/bin/env python3
"""
Generate COPIL Presentation PPT from Notion Database
Automatically creates professional PowerPoint presentations for project governance

Usage:
    python scripts/generate-copil-ppt.py
    python scripts/generate-copil-ppt.py --output my_presentation.pptx
    python scripts/generate-copil-ppt.py --date 2026-04-10
"""

import os
import sys
import json
import argparse
from datetime import datetime
from pathlib import Path

# Third-party imports (install with: pip install python-pptx notion-client python-dateutil pillow matplotlib)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from notion_client import Client

# Load environment variables
from dotenv import load_dotenv

load_dotenv(".env.notion")

NOTION_API_KEY = os.getenv("NOTION_API_KEY")
NOTION_DATABASE_ID = os.getenv("NOTION_DATABASE_ID")

if not NOTION_API_KEY or not NOTION_DATABASE_ID:
    print("❌ Error: NOTION_API_KEY or NOTION_DATABASE_ID not configured")
    print("Set them in .env.notion file")
    sys.exit(1)

# Initialize Notion client
notion = Client(auth=NOTION_API_KEY)

# Color scheme (Banque Marocaine Professional)
COLORS = {
    "primary": RGBColor(0, 51, 102),  # Dark blue
    "secondary": RGBColor(0, 102, 153),  # Medium blue
    "accent": RGBColor(255, 153, 0),  # Orange
    "green": RGBColor(76, 175, 80),  # Green
    "orange": RGBColor(255, 152, 0),  # Orange
    "red": RGBColor(244, 67, 54),  # Red
    "text": RGBColor(33, 33, 33),  # Dark gray
    "light_text": RGBColor(117, 117, 117),  # Light gray
    "white": RGBColor(255, 255, 255),  # White
}


def fetch_tracking_data():
    """Fetch all tracking items from Notion database"""
    print("📥 Fetching data from Notion...")

    try:
        response = notion.databases.query(database_id=NOTION_DATABASE_ID)
        items = []

        for page in response["results"]:
            props = page["properties"]
            item = {
                "id": page["id"],
                "name": (
                    props["Name"]["title"][0]["plain_text"]
                    if props["Name"]["title"]
                    else ""
                ),
                "bloc": props["Bloc"]["select"]["name"]
                if props["Bloc"]["select"]
                else "",
                "statut": props["Statut"]["select"]["name"]
                if props["Statut"]["select"]
                else "",
                "completion": props["Complétion %"]["number"] or 0,
                "description": (
                    props["Description"]["rich_text"][0]["plain_text"]
                    if props["Description"]["rich_text"]
                    else ""
                ),
                "notes": (
                    props["Notes"]["rich_text"][0]["plain_text"]
                    if props["Notes"]["rich_text"]
                    else ""
                ),
            }
            items.append(item)

        print(f"✅ Fetched {len(items)} tracking items")
        return items

    except Exception as e:
        print(f"❌ Error fetching from Notion: {e}")
        sys.exit(1)


def calculate_statistics(items):
    """Calculate KPIs and statistics from tracking items"""
    print("📊 Calculating statistics...")

    stats = {
        "total": len(items),
        "integrated": len([i for i in items if i["statut"] == "INTÉGRÉ"]),
        "in_progress": len([i for i in items if i["statut"] == "EN COURS"]),
        "planned": len([i for i in items if i["statut"] == "PLANIFIÉ"]),
        "blocked": len([i for i in items if i["statut"] == "BLOQUÉ"]),
        "average_completion": sum(i["completion"] for i in items) / len(items)
        if items
        else 0,
        "by_bloc": {},
        "risks": [i for i in items if i["statut"] == "BLOQUÉ" or i["completion"] < 50],
    }

    # Group by bloc
    for item in items:
        bloc = item["bloc"]
        if bloc not in stats["by_bloc"]:
            stats["by_bloc"][bloc] = {
                "total": 0,
                "integrated": 0,
                "in_progress": 0,
                "completion_sum": 0,
                "items": [],
            }

        stats["by_bloc"][bloc]["total"] += 1
        if item["statut"] == "INTÉGRÉ":
            stats["by_bloc"][bloc]["integrated"] += 1
        if item["statut"] == "EN COURS":
            stats["by_bloc"][bloc]["in_progress"] += 1
        stats["by_bloc"][bloc]["completion_sum"] += item["completion"]
        stats["by_bloc"][bloc]["items"].append(item)

    # Calculate average completion per bloc
    for bloc in stats["by_bloc"]:
        stats["by_bloc"][bloc]["completion"] = (
            stats["by_bloc"][bloc]["completion_sum"] / stats["by_bloc"][bloc]["total"]
        )

    return stats


def create_presentation():
    """Create base PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle, date_str):
    """Add title/cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["primary"]

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Date: {date_str}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER


def add_kpi_slide(prs, stats):
    """Add KPI/Executive Summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📊 EXECUTIVE SUMMARY"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # KPI Cards
    kpi_data = [
        ("INTÉGRÉ", f"{stats['integrated']}/{stats['total']}", COLORS["green"]),
        ("EN COURS", f"{stats['in_progress']}/{stats['total']}", COLORS["orange"]),
        ("PLANIFIÉ", f"{stats['planned']}/{stats['total']}", COLORS["secondary"]),
        ("BLOQUÉ", f"{stats['blocked']}/{stats['total']}", COLORS["red"]),
    ]

    y_position = 1.5
    for label, value, color in kpi_data:
        # Card background
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5),
            Inches(y_position),
            Inches(9),
            Inches(1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(1), Inches(y_position + 0.1), Inches(4), Inches(0.4)
        )
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(5), Inches(y_position + 0.1), Inches(4), Inches(0.4)
        )
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.RIGHT

        y_position += 1.2

    # Average completion
    avg_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    avg_frame = avg_box.text_frame
    p = avg_frame.paragraphs[0]
    p.text = f"⏱️  Complétion Moyenne: {stats['average_completion']:.1f}%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]


def add_completion_chart_slide(prs, stats):
    """Add completion by bloc bar chart slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📈 COMPLÉTION PAR BLOC"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # Add bars for each bloc
    y_position = 1.2
    for bloc, data in sorted(
        stats["by_bloc"].items(),
        key=lambda x: x[1]["completion"],
        reverse=True,
    ):
        completion = data["completion"]

        # Bloc label
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_position), Inches(2.5), Inches(0.35)
        )
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = bloc
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]

        # Progress bar background
        bar_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(3.2),
            Inches(y_position + 0.05),
            Inches(6),
            Inches(0.25),
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(200, 200, 200)
        bar_bg.line.color.rgb = RGBColor(150, 150, 150)

        # Progress bar fill
        bar_width = (completion / 100) * 6
        bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(3.2),
            Inches(y_position + 0.05),
            Inches(bar_width),
            Inches(0.25),
        )
        bar.fill.solid()
        # Color based on completion
        if completion >= 90:
            bar.fill.fore_color.rgb = COLORS["green"]
        elif completion >= 70:
            bar.fill.fore_color.rgb = COLORS["accent"]
        else:
            bar.fill.fore_color.rgb = COLORS["orange"]
        bar.line.color.rgb = bar.fill.fore_color.rgb

        # Percentage
        pct_box = slide.shapes.add_textbox(
            Inches(9.3), Inches(y_position), Inches(0.5), Inches(0.35)
        )
        pct_frame = pct_box.text_frame
        p = pct_frame.paragraphs[0]
        p.text = f"{completion:.0f}%"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.RIGHT

        y_position += 0.45


def add_gauge_slide(prs, stats):
    """Add overall completion gauge slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "🎯 COMPLÉTION GLOBALE DU PROJET"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    # Large percentage display
    pct_box = slide.shapes.add_textbox(
        Inches(2), Inches(2.5), Inches(6), Inches(1.5)
    )
    pct_frame = pct_box.text_frame
    p = pct_frame.paragraphs[0]
    p.text = f"{stats['average_completion']:.0f}%"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER

    # Progress bar
    bar_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5),
        Inches(4.2),
        Inches(7),
        Inches(0.5),
    )
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(200, 200, 200)
    bar_bg.line.color.rgb = RGBColor(150, 150, 150)

    # Progress bar fill
    bar_width = (stats["average_completion"] / 100) * 7
    bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5),
        Inches(4.2),
        Inches(bar_width),
        Inches(0.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["green"]
    bar.line.color.rgb = COLORS["green"]

    # Status text
    status_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    status_frame = status_box.text_frame
    status_frame.word_wrap = True
    p = status_frame.paragraphs[0]
    completion = stats["average_completion"]
    if completion >= 90:
        status = "🟢 Excellent - Quasi-finalisé"
    elif completion >= 70:
        status = "🟡 Bon - Sur la bonne voie"
    else:
        status = "🔴 À améliorer - Effort nécessaire"
    p.text = status
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER


def add_risks_slide(prs, risks):
    """Add risks and blockers slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "⚠️  RISQUES & BLOCAGES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    if not risks:
        # No risks
        no_risk_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        no_risk_frame = no_risk_box.text_frame
        p = no_risk_frame.paragraphs[0]
        p.text = "✅ Aucun risque ou blocage identifié"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS["green"]
        p.alignment = PP_ALIGN.CENTER
    else:
        # Display risks
        y_position = 1.2
        for risk in risks[:5]:  # Max 5 risks per slide
            risk_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_position), Inches(8.4), Inches(0.9)
            )
            risk_frame = risk_box.text_frame
            risk_frame.word_wrap = True

            # Risk item with icon
            p = risk_frame.paragraphs[0]
            p.text = f"🟡 {risk['name']}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS["orange"]

            # Details
            p = risk_frame.add_paragraph()
            p.text = f"  Status: {risk['statut']} | Complétion: {risk['completion']}%"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["light_text"]
            p.level = 0

            y_position += 1.1


def add_next_steps_slide(prs):
    """Add next steps/recommendations slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "📋 PROCHAINES ÉTAPES"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    steps = [
        "✅ Phase 9: Mobile Responsive Design (3-5 jours)",
        "✅ Phase 10: E2E Tests avec Cypress (5-7 jours)",
        "✅ Phase 11: Configuration Email & Webhooks (2-3 jours)",
        "✅ Phase 12: Déploiement & Go-Live (2-3 jours)",
    ]

    y_position = 1.5
    for step in steps:
        step_box = slide.shapes.add_textbox(
            Inches(1), Inches(y_position), Inches(8), Inches(0.8)
        )
        step_frame = step_box.text_frame
        step_frame.word_wrap = True
        p = step_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(12)

        y_position += 0.9


def main():
    """Main execution"""
    parser = argparse.ArgumentParser(description="Generate COPIL Presentation PPT")
    parser.add_argument(
        "--output",
        default="COPIL_Presentation.pptx",
        help="Output filename (default: COPIL_Presentation.pptx)",
    )
    parser.add_argument(
        "--date",
        default=datetime.now().strftime("%Y-%m-%d"),
        help="Report date (default: today)",
    )
    args = parser.parse_args()

    print("\n🚀 PF Scoring - COPIL Presentation Generator")
    print("=" * 50)

    # Fetch data
    items = fetch_tracking_data()
    stats = calculate_statistics(items)

    print(f"✅ Statistics calculated:")
    print(f"   - Total items: {stats['total']}")
    print(f"   - Integrated: {stats['integrated']}")
    print(f"   - In progress: {stats['in_progress']}")
    print(f"   - Average completion: {stats['average_completion']:.1f}%")

    # Create presentation
    print("\n📝 Creating presentation...")
    prs = create_presentation()

    # Add slides
    add_title_slide(
        prs,
        "PF SCORING V7++",
        "Rapport de Suivi - COPIL",
        args.date,
    )
    add_kpi_slide(prs, stats)
    add_completion_chart_slide(prs, stats)
    add_gauge_slide(prs, stats)
    add_risks_slide(prs, stats["risks"])
    add_next_steps_slide(prs)

    # Save presentation
    output_path = Path(args.output)
    prs.save(str(output_path))

    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_path.absolute()}")
    print(f"📊 Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
